#!/usr/bin/env python3
"""
PowerPoint Translation Tool using Claude CLI
Reads input.pptx, translates text content using Claude CLI, and outputs output.pptx
"""

import os
import sys
import argparse
import subprocess
import tempfile
import shutil
import json
import time
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx is not installed.")
    print("Please install it using one of these methods:")
    print("1. pip install python-pptx")
    print("2. Create a virtual environment:")
    print("   python3 -m venv venv")
    print("   source venv/bin/activate")
    print("   pip install python-pptx")
    print("3. Or use brew: brew install python-pptx")
    sys.exit(1)


def get_text_hash(text):
    """Get hash of text for caching"""
    import hashlib
    return hashlib.md5(text.encode('utf-8')).hexdigest()

def extract_text_from_shape(shape, slide_num, shape_idx, shape_type_name="unknown"):
    """Extract text from a shape and return structured data"""
    texts = []
    
    # Handle text frames (most common text containers)
    if shape.has_text_frame:
        text_content = shape.text_frame.text
        if text_content.strip():
            texts.append({
                "slide_num": slide_num,
                "shape_idx": shape_idx,
                "shape_type": shape_type_name,
                "text_type": "text_frame",
                "original_text": text_content,
                "translated": False,
                "translated_text": ""
            })
    
    # Handle tables
    if shape.has_table:
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.text_frame.text.strip():
                    texts.append({
                        "slide_num": slide_num,
                        "shape_idx": shape_idx,
                        "shape_type": shape_type_name,
                        "text_type": f"table_cell_{row_idx}_{col_idx}",
                        "original_text": cell.text_frame.text,
                        "translated": False,
                        "translated_text": ""
                    })
    
    # Handle charts (if they have text elements)
    if hasattr(shape, 'has_chart') and shape.has_chart:
        try:
            chart = shape.chart
            # Chart title
            if hasattr(chart, 'chart_title') and chart.chart_title.has_text_frame:
                title_text = chart.chart_title.text_frame.text
                if title_text.strip():
                    texts.append({
                        "slide_num": slide_num,
                        "shape_idx": shape_idx,
                        "shape_type": shape_type_name,
                        "text_type": "chart_title",
                        "original_text": title_text,
                        "translated": False,
                        "translated_text": ""
                    })
        except Exception as e:
            print(f"        Warning: Could not extract chart text: {e}")
    
    # Handle grouped shapes recursively
    if hasattr(shape, 'shapes'):  # This is a group shape
        for sub_shape_idx, sub_shape in enumerate(shape.shapes):
            sub_shape_type = type(sub_shape).__name__
            sub_texts = extract_text_from_shape(sub_shape, slide_num, f"{shape_idx}_{sub_shape_idx}", f"group_{sub_shape_type}")
            texts.extend(sub_texts)
    
    # Only try to extract text from shapes that don't have text_frame to avoid duplicates
    # This is a fallback for shapes that might have text but aren't covered above
    if not shape.has_text_frame:
        try:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append({
                    "slide_num": slide_num,
                    "shape_idx": shape_idx,
                    "shape_type": shape_type_name,
                    "text_type": "shape_text",
                    "original_text": shape.text,
                    "translated": False,
                    "translated_text": ""
                })
        except Exception:
            # This is expected for shapes without text property
            pass
    
    return texts

def extract_text_to_json(input_file, output_json_file):
    """Extract all text from PowerPoint presentation and save to JSON"""
    print(f"Extracting text from presentation: {input_file}")
    
    # Check if output file already exists
    if os.path.exists(output_json_file):
        print(f"JSON file already exists: {output_json_file}")
        print("Skipping text extraction (file already exists)")
        return True
    
    # Load presentation
    try:
        prs = Presentation(input_file)
    except Exception as e:
        print(f"Error loading presentation: {e}")
        return False
    
    total_slides = len(prs.slides)
    print(f"Processing {total_slides} slides for text extraction...")
    
    all_texts = []
    
    # Process each slide
    for slide_num, slide in enumerate(prs.slides):
        print(f"  [{slide_num + 1}/{total_slides}] Extracting text from slide {slide_num + 1}")
        
        # Process all shapes in the slide
        for shape_idx, shape in enumerate(slide.shapes):
            shape_type_name = type(shape).__name__
            print(f"    Processing shape {shape_idx}: {shape_type_name}")
            
            try:
                texts = extract_text_from_shape(shape, slide_num + 1, shape_idx, shape_type_name)
                if texts:
                    all_texts.extend(texts)
                    print(f"      Found {len(texts)} text elements")
                else:
                    print(f"      No text found in shape")
            except Exception as e:
                print(f"    Warning: Error extracting text from shape: {e}")
        
        # Process placeholders that might not be in shapes
        for placeholder_idx, placeholder in enumerate(slide.placeholders):
            placeholder_type_name = type(placeholder).__name__
            print(f"    Processing placeholder {placeholder_idx}: {placeholder_type_name}")
            
            try:
                texts = extract_text_from_shape(placeholder, slide_num + 1, f"placeholder_{placeholder_idx}", f"placeholder_{placeholder_type_name}")
                if texts:
                    all_texts.extend(texts)
                    print(f"      Found {len(texts)} text elements in placeholder")
                else:
                    print(f"      No text found in placeholder")
            except Exception as e:
                print(f"    Warning: Error extracting text from placeholder: {e}")
    
    # Save to JSON file
    try:
        with open(output_json_file, 'w', encoding='utf-8') as f:
            json.dump(all_texts, f, ensure_ascii=False, indent=2)
        
        print(f"✓ Successfully extracted {len(all_texts)} text elements to: {output_json_file}")
        return True
    except Exception as e:
        print(f"Error saving JSON file: {e}")
        return False

def translate_json_texts(json_file, output_lang, custom_prompt=None, max_retries=3):
    """Translate texts in JSON file and update translation status"""
    print(f"Translating texts from JSON file: {json_file}")
    
    # Load JSON file
    if not os.path.exists(json_file):
        print(f"Error: JSON file not found: {json_file}")
        return False
    
    try:
        with open(json_file, 'r', encoding='utf-8') as f:
            texts_data = json.load(f)
    except Exception as e:
        print(f"Error loading JSON file: {e}")
        return False
    
    if not texts_data:
        print("No text data found in JSON file")
        return True
    
    total_texts = len(texts_data)
    translated_count = 0
    skipped_count = 0
    failed_count = 0
    
    print(f"Processing {total_texts} text elements...")
    
    # Build a cache of existing translations for quick lookup
    print("Building translation cache from existing translations...")
    translation_cache = {}
    for item in texts_data:
        if item.get("translated", False) and item.get("original_text") and item.get("translated_text"):
            original = item.get("original_text", "").strip()
            translated = item.get("translated_text", "").strip()
            if original and translated:
                translation_cache[original] = translated
    
    print(f"Found {len(translation_cache)} existing translations in cache")
    
    for idx, text_item in enumerate(texts_data):
        print(f"  [{idx + 1}/{total_texts}] Processing text item...")
        
        # Skip if already translated
        if text_item.get("translated", False):
            print(f"    Skipping: Already translated")
            skipped_count += 1
            continue
        
        original_text = text_item.get("original_text", "")
        if not original_text or not original_text.strip():
            print(f"    Skipping: Empty text")
            skipped_count += 1
            continue
        
        print(f"    Original: '{original_text[:50]}...'")
        
        # Check if we have this translation in cache
        if original_text.strip() in translation_cache:
            cached_translation = translation_cache[original_text.strip()]
            print(f"    Found in cache: '{cached_translation[:50]}...'")
            text_item["translated"] = True
            text_item["translated_text"] = cached_translation
            translated_count += 1
            
            # Save progress after using cache
            try:
                with open(json_file, 'w', encoding='utf-8') as f:
                    json.dump(texts_data, f, ensure_ascii=False, indent=2)
            except Exception as e:
                print(f"    Warning: Could not save progress: {e}")
            continue
        
        # Attempt translation with retries
        success = False
        for attempt in range(max_retries):
            if attempt > 0:
                print(f"    Retry attempt {attempt + 1}/{max_retries}")
            
            try:
                translation_result = translate_text_with_claude(original_text, output_lang, custom_prompt, 1)
                
                if translation_result and not translation_result.get("is_error", False):
                    translated_text = translation_result.get("text", "")
                    if translated_text and translated_text.strip():
                        # Update the text item with translation
                        text_item["translated"] = True
                        text_item["translated_text"] = translated_text
                        
                        # Add to cache for future use
                        translation_cache[original_text.strip()] = translated_text.strip()
                        
                        print(f"    Translated: '{translated_text[:50]}...'")
                        success = True
                        translated_count += 1
                        break
                    else:
                        print(f"    Attempt {attempt + 1}: Empty translation result")
                else:
                    print(f"    Attempt {attempt + 1}: Translation failed")
                    
            except Exception as e:
                print(f"    Attempt {attempt + 1}: Exception: {e}")
        
        # If all retries failed, mark as failed
        if not success:
            print(f"    Failed: All {max_retries} attempts failed")
            text_item["translated"] = False
            text_item["translated_text"] = ""
            failed_count += 1
        
        # Save progress after each translation
        try:
            with open(json_file, 'w', encoding='utf-8') as f:
                json.dump(texts_data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"    Warning: Could not save progress: {e}")
    
    # Final save
    try:
        with open(json_file, 'w', encoding='utf-8') as f:
            json.dump(texts_data, f, ensure_ascii=False, indent=2)
        print(f"✓ Successfully saved updated JSON file")
    except Exception as e:
        print(f"Error saving final JSON file: {e}")
        return False
    
    print(f"\n=== Translation Summary ===")
    print(f"Total texts: {total_texts}")
    print(f"Successfully translated: {translated_count}")
    print(f"Skipped (already translated): {skipped_count}")
    print(f"Failed: {failed_count}")
    
    # Check if all translations are complete
    all_complete = (translated_count + skipped_count == total_texts) and (failed_count == 0)
    
    if all_complete:
        print(f"\n✅ All translations completed successfully!")
        print(f"Ready to apply translations to PowerPoint file.")
    elif failed_count > 0:
        print(f"\n⚠️  {failed_count} translations failed.")
        print(f"You may want to retry failed translations or proceed with partial results.")
    else:
        print(f"\n📝 Translation in progress...")
    
    return True

def apply_translations_to_ppt(input_ppt_file, json_file, output_ppt_file):
    """Apply translations from JSON file to PowerPoint file"""
    print(f"Applying translations from JSON to PowerPoint...")
    print(f"Input PPT: {input_ppt_file}")
    print(f"JSON file: {json_file}")
    print(f"Output PPT: {output_ppt_file}")
    
    # Load JSON file
    if not os.path.exists(json_file):
        print(f"Error: JSON file not found: {json_file}")
        return False
    
    try:
        with open(json_file, 'r', encoding='utf-8') as f:
            texts_data = json.load(f)
    except Exception as e:
        print(f"Error loading JSON file: {e}")
        return False
    
    # Load PowerPoint presentation
    try:
        prs = Presentation(input_ppt_file)
    except Exception as e:
        print(f"Error loading PowerPoint file: {e}")
        return False
    
    if not texts_data:
        print("No text data found in JSON file")
        return False
    
    total_texts = len(texts_data)
    applied_count = 0
    skipped_count = 0
    failed_count = 0
    
    print(f"Processing {total_texts} text elements...")
    
    # Group texts by slide for easier processing
    texts_by_slide = {}
    for text_item in texts_data:
        slide_num = text_item.get("slide_num", 1)
        if slide_num not in texts_by_slide:
            texts_by_slide[slide_num] = []
        texts_by_slide[slide_num].append(text_item)
    
    # Process each slide
    for slide_num, slide_texts in texts_by_slide.items():
        print(f"  Processing slide {slide_num} ({len(slide_texts)} texts)...")
        
        # Get slide (convert to 0-based index)
        slide_idx = slide_num - 1
        if slide_idx >= len(prs.slides):
            print(f"    Warning: Slide {slide_num} not found in presentation")
            failed_count += len(slide_texts)
            continue
        
        slide = prs.slides[slide_idx]
        
        # Process each text item for this slide
        for text_item in slide_texts:
            shape_idx = text_item.get("shape_idx")
            text_type = text_item.get("text_type", "")
            original_text = text_item.get("original_text", "")
            translated = text_item.get("translated", False)
            translated_text = text_item.get("translated_text", "")
            
            # Skip if not translated
            if not translated or not translated_text:
                print(f"    Skipping: Not translated - {original_text[:30]}...")
                skipped_count += 1
                continue
            
            # Try to find and replace the text
            try:
                success = replace_text_in_slide(slide, shape_idx, text_type, original_text, translated_text)
                if success:
                    print(f"    Applied: {original_text[:30]}... → {translated_text[:30]}...")
                    applied_count += 1
                else:
                    print(f"    Failed: Could not locate text - {original_text[:30]}...")
                    failed_count += 1
            except Exception as e:
                print(f"    Error: {e}")
                failed_count += 1
    
    # Save the updated presentation
    try:
        prs.save(output_ppt_file)
        print(f"✓ Successfully saved translated PowerPoint file: {output_ppt_file}")
    except Exception as e:
        print(f"Error saving PowerPoint file: {e}")
        return False
    
    print(f"\n=== Translation Application Summary ===")
    print(f"Total texts: {total_texts}")
    print(f"Successfully applied: {applied_count}")
    print(f"Skipped (not translated): {skipped_count}")
    print(f"Failed: {failed_count}")
    
    if applied_count > 0:
        print(f"\n✅ Translations successfully applied to PowerPoint file!")
        print(f"Output saved to: {output_ppt_file}")
    
    return True

def replace_text_in_slide(slide, shape_idx, text_type, original_text, translated_text):
    """Replace text in a specific shape within a slide"""
    
    # Handle different shape index formats
    if isinstance(shape_idx, str):
        if shape_idx.startswith("placeholder_"):
            # Handle placeholder
            placeholder_idx = int(shape_idx.split("_")[1])
            if placeholder_idx < len(slide.placeholders):
                shape = slide.placeholders[placeholder_idx]
                return replace_text_in_shape(shape, text_type, original_text, translated_text)
        elif "_" in shape_idx:
            # Handle grouped shape (e.g., "2_1", "2_2_0")
            parts = shape_idx.split("_")
            try:
                main_idx = int(parts[0])
                if main_idx < len(slide.shapes):
                    shape = slide.shapes[main_idx]
                    # Navigate to sub-shape
                    for sub_idx_str in parts[1:]:
                        sub_idx = int(sub_idx_str)
                        if hasattr(shape, 'shapes') and sub_idx < len(shape.shapes):
                            shape = shape.shapes[sub_idx]
                        else:
                            return False
                    return replace_text_in_shape(shape, text_type, original_text, translated_text)
            except (ValueError, IndexError):
                return False
    else:
        # Handle regular shape index
        if shape_idx < len(slide.shapes):
            shape = slide.shapes[shape_idx]
            return replace_text_in_shape(shape, text_type, original_text, translated_text)
    
    return False

def replace_text_in_shape(shape, text_type, original_text, translated_text):
    """Replace text in a specific shape"""
    try:
        if text_type == "text_frame" and shape.has_text_frame:
            # Replace text in text frame
            if shape.text_frame.text == original_text:
                shape.text_frame.clear()
                paragraph = shape.text_frame.paragraphs[0]
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.text = translated_text
                return True
        
        elif text_type.startswith("table_cell_") and shape.has_table:
            # Replace text in table cell
            parts = text_type.split("_")
            if len(parts) >= 4:
                row_idx = int(parts[2])
                col_idx = int(parts[3])
                table = shape.table
                if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                    cell = table.rows[row_idx].cells[col_idx]
                    if cell.text_frame.text == original_text:
                        cell.text_frame.clear()
                        paragraph = cell.text_frame.paragraphs[0]
                        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        run.text = translated_text
                        return True
        
        elif text_type == "chart_title" and hasattr(shape, 'has_chart') and shape.has_chart:
            # Replace text in chart title
            chart = shape.chart
            if hasattr(chart, 'chart_title') and chart.chart_title.has_text_frame:
                if chart.chart_title.text_frame.text == original_text:
                    chart.chart_title.text_frame.clear()
                    paragraph = chart.chart_title.text_frame.paragraphs[0]
                    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                    run.text = translated_text
                    return True
        
        elif text_type == "shape_text" and hasattr(shape, 'text'):
            # Replace text in shape text property
            if shape.text == original_text:
                shape.text = translated_text
                return True
        
    except Exception as e:
        print(f"      Error replacing text: {e}")
        return False
    
    return False


def get_language_name(lang_code):
    """Convert language code to full name"""
    lang_map = {
        'zh': 'Chinese',
        'en': 'English',
        'ja': 'Japanese',
        'ko': 'Korean',
        'fr': 'French',
        'de': 'German',
        'es': 'Spanish',
        'it': 'Italian',
        'pt': 'Portuguese',
        'ru': 'Russian',
        'ar': 'Arabic',
        'hi': 'Hindi',
        'th': 'Thai',
        'vi': 'Vietnamese'
    }
    return lang_map.get(lang_code.lower(), lang_code)

def create_translation_prompt(output_lang, custom_prompt=None):
    """Create translation prompt for PowerPoint text"""
    lang_name = get_language_name(output_lang)
    
    base_prompt = f"""请翻译PowerPoint文本为 {lang_name}. 
CRITICAL: 你必须严格按照指定格式输出，否则会导致处理失败！

MANDATORY FORMAT REQUIREMENTS:
1. 第一行必须是：<!-- TRANSLATION_START -->
2. 中间行：翻译后的文本内容（可以是多行）
3. 最后一行必须是：<!-- TRANSLATION_END -->
4. 绝对不能有任何其他内容、说明、注释或解释
5. 如果无法翻译，请输出：<!-- TRANSLATION_START -->\n原文内容\n<!-- TRANSLATION_END -->

TRANSLATION REQUIREMENTS:
1. 仅翻译文字内容，保持简洁
2. 保持术语一致性
3. 适合演示文稿的简洁表达
4. 删除不必要的字符和标记
5. 使用 Sonnet 4 模型完成，保证格式和语义准确翻译内容自然流畅
6. 表达清晰简洁，适合PPT展示

WARNING: 如果不按照上述格式输出，翻译将被视为失败并重试！"""    
    if custom_prompt:
        base_prompt += f"\n\nADDITIONAL INSTRUCTIONS:\n{custom_prompt}"
    
    base_prompt += "\n\n待翻译文本:"
    
    return base_prompt

def translate_text_with_claude(text, output_lang, custom_prompt=None, max_retries=3):
    """Translate text using Claude CLI with retry mechanism"""
    if not text or not text.strip():
        return {"text": text, "is_error": False}
    
    # Check if translation failed and return original with error flag
    def return_with_error(original_text, error_reason=""):
        print(f"        Translation failed: {error_reason}")
        return {"text": original_text, "is_error": True}
    
    # Try translation with retries
    for attempt in range(max_retries):
        if attempt > 0:
            print(f"        Retry attempt {attempt + 1}/{max_retries}")
        
        # Create translation prompt
        prompt = create_translation_prompt(output_lang, custom_prompt)
        
        try:
            # Prepare the full input text
            full_input = f"{prompt}\n\n{text}"
            
            # Use Claude CLI
            claude_command = ['claude']
            
            # Run Claude CLI command
            # Set environment variables to skip interactive prompts
            env = os.environ.copy()
            env['CLAUDE_NON_INTERACTIVE'] = '1'
            env['CLAUDE_AUTO_APPROVE'] = '1'
            
            result = subprocess.run(
                claude_command,
                input=full_input,
                capture_output=True,
                text=True,
                timeout=120,  # 2 minute timeout
                encoding='utf-8',
                env=env
            )
            
            if result.returncode == 0:
                translated_text = result.stdout.strip()
                
                # Extract content between <!-- TRANSLATION_START --> and <!-- TRANSLATION_END --> markers
                start_marker = '<!-- TRANSLATION_START -->'
                end_marker = '<!-- TRANSLATION_END -->'
                
                start_idx = translated_text.find(start_marker)
                end_idx = translated_text.find(end_marker)
                
                if start_idx != -1 and end_idx != -1:
                    # Extract the content between markers
                    content_start = start_idx + len(start_marker)
                    extracted_content = translated_text[content_start:end_idx].strip()
                    
                    # Validate extracted content
                    if extracted_content and len(extracted_content.strip()) > 0:
                        return {"text": extracted_content, "is_error": False}
                    else:
                        print(f"        Attempt {attempt + 1}: Empty content between markers")
                        continue  # Retry
                else:
                    # If required markers not found, this is an error
                    print(f"        Attempt {attempt + 1}: Required markers <!-- TRANSLATION_START --> and <!-- TRANSLATION_END --> not found")
                    print(f"        Raw output: {translated_text[:200]}...")
                    continue  # Retry
            else:
                error_msg = result.stderr.strip() if result.stderr else "No error message"
                print(f"        Attempt {attempt + 1}: CLI error (code {result.returncode}): {error_msg}")
                continue  # Retry
                
        except subprocess.TimeoutExpired:
            print(f"        Attempt {attempt + 1}: Translation timeout (2 minutes)")
            continue  # Retry
        except FileNotFoundError:
            return return_with_error(text, "'claude' command not found")
        except Exception as e:
            print(f"        Attempt {attempt + 1}: Exception: {e}")
            continue  # Retry
    
    # All retries failed
    return return_with_error(text, f"Translation failed after {max_retries} attempts")



def check_claude_cli():
    """Check if Claude CLI is available"""
    try:
        result = subprocess.run(['claude', '--version'], 
                              capture_output=True, text=True, timeout=10)
        if result.returncode == 0:
            version_info = result.stdout.strip().split('\n')[-1]
            print(f"Claude CLI available: {version_info}")
            return True
        else:
            print(f"Claude CLI check failed with code {result.returncode}")
            return False
    except subprocess.TimeoutExpired:
        print("Error: Claude CLI version check timed out")
        return False
    except FileNotFoundError:
        print("Error: 'claude' command not found")
        print("Please ensure Claude CLI is installed and in your PATH")
        return False
    except Exception as e:
        print(f"Error checking Claude CLI: {e}")
        return False

def parse_arguments():
    """Parse command line arguments"""
    parser = argparse.ArgumentParser(
        description="PowerPoint Translation Tool using Claude CLI"
    )
    
    parser.add_argument(
        'input_file',
        default='input.pptx',
        nargs='?',
        help="Input PowerPoint file (default: input.pptx)"
    )
    
    parser.add_argument(
        '-o', '--output',
        default='output.pptx',
        help="Output PowerPoint file (default: output.pptx)"
    )
    
    parser.add_argument(
        '--olang',
        required=False,
        help="Target language code (e.g., zh, en, ja, ko, fr, de, es, it, pt, ru, ar, hi, th, vi)"
    )
    
    parser.add_argument(
        '-p', '--prompt',
        default=None,
        help="Additional custom prompt to add to the translation instructions"
    )
    
    
    parser.add_argument(
        '--extract-text',
        action='store_true',
        help="Extract text from PowerPoint to JSON file (skips translation)"
    )
    
    parser.add_argument(
        '--json-output',
        default=None,
        help="Output JSON file for text extraction (default: <input_basename>_texts.json)"
    )
    
    parser.add_argument(
        '--translate-json',
        action='store_true',
        help="Translate texts in existing JSON file"
    )
    
    parser.add_argument(
        '--json-input',
        default=None,
        help="Input JSON file for translation (default: <input_basename>_texts.json)"
    )
    
    parser.add_argument(
        '--apply-translations',
        action='store_true',
        help="Apply translations from JSON file back to PowerPoint"
    )
    
    return parser.parse_args()

def main():
    """Main function"""
    print("=== PowerPoint Translation Tool ===")
    
    # Parse arguments
    args = parse_arguments()
    
    # Check input file (only for non-JSON translation modes)
    if not args.translate_json and not os.path.exists(args.input_file):
        print(f"Error: Input file not found: {args.input_file}")
        sys.exit(1)
    
    # Handle text extraction mode
    if args.extract_text:
        print("Mode: Text Extraction")
        
        # Create temp directory
        input_basename = os.path.splitext(os.path.basename(args.input_file))[0]
        temp_dir = f"{input_basename}_temp"
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
        
        # Determine JSON output file
        if args.json_output:
            json_output = args.json_output
        else:
            json_output = os.path.join(temp_dir, f"{input_basename}_texts.json")
        
        print(f"Input file: {args.input_file}")
        print(f"Temp directory: {temp_dir}")
        print(f"JSON output: {json_output}")
        
        # Extract text to JSON
        success = extract_text_to_json(args.input_file, json_output)
        
        if success:
            print("\n=== Text Extraction Complete ===")
            print(f"Text data saved to: {json_output}")
        else:
            print("\n=== Text Extraction Failed ===")
            sys.exit(1)
        
        return
    
    # Handle JSON translation mode
    if args.translate_json:
        print("Mode: JSON Translation")
        
        # Check required arguments for translation
        if not args.olang:
            print("Error: --olang is required for JSON translation mode")
            sys.exit(1)
        
        # Check Claude CLI availability
        if not check_claude_cli():
            sys.exit(1)
        
        # Create temp directory with language identifier
        input_basename = os.path.splitext(os.path.basename(args.input_file))[0]
        temp_dir = f"{input_basename}_temp_{args.olang}"
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
        
        # Determine JSON input file
        if args.json_input:
            json_input = args.json_input
        else:
            json_input = os.path.join(temp_dir, f"{input_basename}_texts.json")
        
        print(f"Temp directory: {temp_dir}")
        print(f"JSON input: {json_input}")
        print(f"Target language: {args.olang}")
        
        if args.prompt:
            print(f"Custom prompt: {args.prompt}")
        
        # Translate JSON texts
        success = translate_json_texts(json_input, args.olang, args.prompt)
        
        if success:
            print("\n=== JSON Translation Complete ===")
            print(f"Updated JSON file: {json_input}")
        else:
            print("\n=== JSON Translation Failed ===")
            sys.exit(1)
        
        return
    
    # Handle apply translations mode
    if args.apply_translations:
        print("Mode: Apply Translations")
        
        # Create temp directory with language identifier if available
        input_basename = os.path.splitext(os.path.basename(args.input_file))[0]
        temp_dir = f"{input_basename}_temp_{args.olang}" if args.olang else f"{input_basename}_temp"
        if not os.path.exists(temp_dir):
            os.makedirs(temp_dir)
        
        # Determine JSON input file
        if args.json_input:
            json_input = args.json_input
        else:
            json_input = os.path.join(temp_dir, f"{input_basename}_texts.json")
        
        print(f"Input PPT: {args.input_file}")
        print(f"Temp directory: {temp_dir}")
        print(f"JSON input: {json_input}")
        print(f"Output PPT: {args.output}")
        
        # Apply translations to PowerPoint
        success = apply_translations_to_ppt(args.input_file, json_input, args.output)
        
        if success:
            print("\n=== Apply Translations Complete ===")
            print(f"Translated PowerPoint saved to: {args.output}")
        else:
            print("\n=== Apply Translations Failed ===")
            sys.exit(1)
        
        return
    
    # Default mode: Full translation (3-step process)
    print("Mode: Full Translation (3-step process)")
    
    # Check required arguments for translation
    if not args.olang:
        print("Error: --olang is required for translation mode")
        sys.exit(1)
    
    # Check Claude CLI availability
    if not check_claude_cli():
        sys.exit(1)
    
    # Create temp directory based on input filename with language identifier
    input_basename = os.path.splitext(os.path.basename(args.input_file))[0]
    temp_dir = f"{input_basename}_temp_{args.olang}"
    
    if not os.path.exists(temp_dir):
        os.makedirs(temp_dir)
    
    json_file = os.path.join(temp_dir, f"{input_basename}_texts.json")
    
    print(f"Input file: {args.input_file}")
    print(f"Output file: {args.output}")
    print(f"Target language: {args.olang}")
    print(f"Temp directory: {temp_dir}")
    print(f"JSON file: {json_file}")
    
    if args.prompt:
        print(f"Custom prompt: {args.prompt}")
    
    # Step 1: Extract text to JSON
    print("\n=== Step 1: Extracting text to JSON ===")
    success = extract_text_to_json(args.input_file, json_file)
    if not success:
        print("=== Step 1 Failed ===")
        sys.exit(1)
    print("=== Step 1 Complete ===")
    
    # Step 2: Translate JSON texts
    print("\n=== Step 2: Translating texts ===")
    success = translate_json_texts(json_file, args.olang, args.prompt)
    if not success:
        print("=== Step 2 Failed ===")
        sys.exit(1)
    print("=== Step 2 Complete ===")
    
    # Step 3: Apply translations to PowerPoint
    print("\n=== Step 3: Applying translations to PowerPoint ===")
    success = apply_translations_to_ppt(args.input_file, json_file, args.output)
    if not success:
        print("=== Step 3 Failed ===")
        sys.exit(1)
    print("=== Step 3 Complete ===")
    
    # Clean up temp directory (optional, keep for debugging)
    # if os.path.exists(temp_dir):
    #     shutil.rmtree(temp_dir)
    
    print("\n=== Full Translation Complete ===")
    print(f"Translated presentation saved to: {args.output}")
    print(f"Intermediate files saved in: {temp_dir}")
    print("You can reuse the JSON file for future translations or modifications.")

if __name__ == "__main__":
    main()